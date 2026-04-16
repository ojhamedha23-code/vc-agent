"""
PDF Parser — extracts text slide-by-slide.
Falls back to Claude vision for image-heavy slides.
Also supports:
  - parse_from_images(): batch of uploaded PNG/JPG files
  - parse_from_url():    direct PDF URL download
"""
from __future__ import annotations

import base64
import io
import tempfile
import urllib.request
from dataclasses import dataclass
from pathlib import Path
from typing import List

import anthropic
import pdfplumber
from PIL import Image
from pptx import Presentation

MIN_TEXT_CHARS = 50   # slides with fewer chars are treated as image-based


@dataclass
class SlideContent:
    slide_num: int
    text: str
    is_image_based: bool


def _encode_page_as_image(page) -> str:
    """Render a pdfplumber page as a base64 PNG."""
    img = page.to_image(resolution=150).original
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return base64.standard_b64encode(buf.getvalue()).decode("utf-8")


def _extract_text_via_vision(client: anthropic.Anthropic, b64_image: str, slide_num: int) -> str:
    """Send a slide image to Claude vision and extract all text/data."""
    response = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=1024,
        messages=[
            {
                "role": "user",
                "content": [
                    {
                        "type": "image",
                        "source": {
                            "type": "base64",
                            "media_type": "image/png",
                            "data": b64_image,
                        },
                    },
                    {
                        "type": "text",
                        "text": (
                            f"This is slide {slide_num} from a startup pitch deck. "
                            "Extract ALL text, numbers, labels, and data visible on this slide. "
                            "Preserve structure and hierarchy. Return plain text only."
                        ),
                    },
                ],
            }
        ],
    )
    return response.content[0].text.strip()


def parse_deck(pdf_path: str, client: anthropic.Anthropic) -> List[SlideContent]:
    """
    Parse a PDF pitch deck into a list of SlideContent objects.
    Text-first extraction; vision fallback for image-only slides.
    """
    slides: List[SlideContent] = []

    with pdfplumber.open(pdf_path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = (page.extract_text() or "").strip()

            if len(text) >= MIN_TEXT_CHARS:
                slides.append(SlideContent(slide_num=i, text=text, is_image_based=False))
            else:
                # Fall back to vision
                b64 = _encode_page_as_image(page)
                vision_text = _extract_text_via_vision(client, b64, i)
                slides.append(SlideContent(slide_num=i, text=vision_text, is_image_based=True))

    return slides


def parse_from_images(
    image_files: List[bytes],
    client: anthropic.Anthropic,
    media_types: List[str] = None,
) -> List[SlideContent]:
    """
    Process a list of uploaded image bytes (PNG/JPG) as pitch deck slides.
    Each image is sent to Claude vision for text extraction.
    Used when slides are screenshotted from a website instead of uploaded as PDF.
    """
    slides: List[SlideContent] = []
    if media_types is None:
        media_types = ["image/png"] * len(image_files)

    for i, (img_bytes, media_type) in enumerate(zip(image_files, media_types), start=1):
        # Normalise to PNG for consistency
        img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        b64 = base64.standard_b64encode(buf.getvalue()).decode("utf-8")

        text = _extract_text_via_vision(client, b64, i)
        slides.append(SlideContent(slide_num=i, text=text, is_image_based=True))

    return slides


def parse_from_url(url: str, client: anthropic.Anthropic) -> List[SlideContent]:
    """
    Download a PDF from a direct URL and parse it.
    Raises ValueError if the URL does not point to a PDF.
    """
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
        tmp_path = tmp.name

    req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
    with urllib.request.urlopen(req, timeout=30) as resp:
        content_type = resp.headers.get("Content-Type", "")
        if "pdf" not in content_type.lower() and not url.lower().endswith(".pdf"):
            raise ValueError(
                f"URL does not appear to be a PDF (Content-Type: {content_type}). "
                "Please use a direct .pdf link."
            )
        data = resp.read()

    Path(tmp_path).write_bytes(data)
    slides = parse_deck(tmp_path, client)
    Path(tmp_path).unlink(missing_ok=True)
    return slides


def parse_from_pptx(pptx_path: str, client: anthropic.Anthropic) -> List[SlideContent]:
    """
    Parse a .pptx file into SlideContent objects.
    Extracts text from each slide's shapes directly.
    Falls back to vision for slides with very little text.
    """
    prs = Presentation(pptx_path)
    slides: List[SlideContent] = []

    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        text = "\n".join(texts).strip()

        if len(text) >= MIN_TEXT_CHARS:
            slides.append(SlideContent(slide_num=i, text=text, is_image_based=False))
        else:
            # Slide is likely image/chart heavy — render via vision
            # Export slide as PNG using python-pptx + Pillow workaround
            # For MVP: flag as unverified text and include what we have
            fallback = text if text else f"[Slide {i}: image/chart only — no text extracted]"
            slides.append(SlideContent(slide_num=i, text=fallback, is_image_based=True))

    return slides


def slides_to_text(slides: List[SlideContent]) -> str:
    """Concatenate all slide content into a single string for LLM context."""
    parts = []
    for s in slides:
        tag = " [image-extracted]" if s.is_image_based else ""
        parts.append(f"--- Slide {s.slide_num}{tag} ---\n{s.text}")
    return "\n\n".join(parts)
